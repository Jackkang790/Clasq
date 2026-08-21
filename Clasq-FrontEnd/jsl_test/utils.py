#src/utils 현재경로


#========================core.py=====================

# =========================================================
# [core.py]
# 통합 코어 모듈 - 파일 관리 시스템의 핵심 기능을 결합
# =========================================================
import os
import re
import shutil
from typing import Dict, Any, List

from .file_pipeline import TextExtractor, FileAnalyzer
from .query_parser import SearchQueryParser
from .db_manager import FileRegistryManager
from .search_engine import SearchEngine


class ClasqCore:
    """
    통합 코어 클래스
    파일 분석, DB 관리, 검색 엔진을 하나로 결합하여
    시스템의 핵심 기능을 제공합니다.
    """

    def __init__(
        self,
        db_path: str = "file_manager.db",
        ollama_url: str = "http://localhost:11434",
        text_model: str = "gemma2:9b",
        vision_model: str = "llava",
        whisper_model: str = "base"
    ):
        """
        코어 시스템 초기화
        
        Args:
            db_path: SQLite DB 파일 경로
            ollama_url: Ollama API URL
            text_model: 텍스트 분석용 모델
            vision_model: 이미지 분석용 모델
            whisper_model: 음성 인식용 모델
        """
        self.db_path = db_path
        
        # DB 관리자 초기화
        self.registry = FileRegistryManager(db_path=db_path)
        
        # 파일 파이프라인 초기화
        self.extractor = TextExtractor(whisper_model_name=whisper_model)
        self.analyzer = FileAnalyzer(
            ollama_url=ollama_url,
            text_model=text_model,
            vision_model=vision_model
        )
        
        # 검색 파서 초기화
        self.query_parser = SearchQueryParser(
            ollama_url=ollama_url,
            model=text_model
        )
        
        # 검색 엔진 초기화
        self.search_engine = SearchEngine(db_path=db_path)

    def _normalize_path(self, path: str) -> str:
        """경로 문자열 정제 (윈도우 경로 깨짐 방어)"""
        if not path:
            return ""
        clean_path = path.replace('￥', '/').replace('\\', '/')
        return os.path.abspath(clean_path)

    def _save_to_db(self, file_path: str, metadata_result: Dict[str, Any]) -> Dict[str, Any]:
        """AI 분석 결과를 DB에 저장"""
        result = self.registry.save_file_result(file_path, metadata_result)
        if not result.get("success"):
            print(f"[DB 저장 오류]: {result.get('message')}")
        elif result.get("is_duplicate"):
            print(f"[중복 파일 감지]: {file_path} -> {result.get('duplicate_of')} 와 내용 동일 "
                  f"(정책: {self.registry.duplicate_policy})")
        return result

    @staticmethod
    def _add_media_filename_tags(file_path: str, result: Dict[str, Any]) -> None:
        """미디어 AI 분석이 실패하거나 태그가 비어 있을 때 파일명 기반 태그를 보완한다."""
        metadata = result.get("metadata")
        if not isinstance(metadata, dict):
            metadata = {}
            result["metadata"] = metadata
        tags = metadata.get("tags")
        valid_tags = [
            tag.strip()
            for tag in tags
            if isinstance(tag, str) and tag.strip()
        ] if isinstance(tags, list) else []
        if valid_tags:
            metadata["tags"] = valid_tags
            return

        extension = os.path.splitext(file_path)[1].lower().lstrip(".")
        media_kind = "오디오" if extension in {"mp3", "wav", "m4a"} else "영상"
        filename = os.path.splitext(os.path.basename(file_path))[0]
        filename_tokens = re.findall(r"[0-9A-Za-z가-힣]+", filename)

        fallback_tags = [media_kind, extension or "미디어"]
        for token in filename_tokens:
            if token.lower() not in {tag.lower() for tag in fallback_tags}:
                fallback_tags.append(token)
            if len(fallback_tags) == 5:
                break

        description = str(metadata.get("description") or "").strip()
        metadata["tags"] = fallback_tags
        metadata["description"] = (
            description or "AI 분석 결과가 없어 파일명과 확장자를 기준으로 태그를 부착했습니다."
        )
        metadata["ai_comment"] = (
            f"태그: {', '.join(f'#{tag}' for tag in fallback_tags)} / "
            f"코멘트: {metadata['description']}"
        )
        result["tag_source"] = "filename_fallback"

    def sync_db_with_disk(self):
        """DB와 실제 디스크 파일 동기화"""
        return self.registry.sync_missing_files()

    # ---------------------------------------------------------
    # [유스케이스 1] 파일 업로드 및 분석 요청 처리
    # ---------------------------------------------------------
    def process_file_upload(self, raw_file_path: str) -> Dict[str, Any]:
        """
        파일 처리 파이프라인
        경로 정제 -> 파일 종류 판별 -> 전처리 -> AI 분석 -> DB 저장
        """
        file_path = self._normalize_path(raw_file_path)

        if not os.path.exists(file_path):
            return {
                "@TYPE": "@ERROR",
                "message": f"파일을 찾을 수 없습니다: {file_path}"
            }

        # A. 이미지 파일 처리
        if self.extractor.is_image_file(file_path):
            img_bytes, status = self.extractor.process_image(file_path)
            if status != "SUCCESS":
                res = self.analyzer._build_fallback_response(
                    {"original_name": os.path.basename(file_path)}, status)
            else:
                res = self.analyzer.analyze_image_bytes(file_path, img_bytes)

        # B. 오디오/비디오 미디어 파일 처리
        elif self.extractor.is_media_file(file_path):
            text, status = self.extractor.process_media(file_path)
            res = self.analyzer.analyze_document_text(file_path, text)
            self._add_media_filename_tags(file_path, res)

        # C. 일반 문서/데이터 파일 처리
        else:
            text, status = self.extractor.extract(file_path)
            res = self.analyzer.analyze_document_text(file_path, text)

        # 분석 결과를 DB에 저장
        db_save = self._save_to_db(file_path, res)
        res["db_save"] = db_save

        return res

    # ---------------------------------------------------------
    # [유스케이스 2] 자연어 검색창 입력문 처리
    # ---------------------------------------------------------
    def process_user_query(self, user_text: str) -> Dict[str, Any]:
        """
        사용자 자연어 입력 처리
        의도 파싱 -> 검색 엔진 전달 -> 결과 반환
        """
        # 1단계: 자연어 의도 파싱
        parse_result = self.query_parser.parse_user_query(user_text)
        
        if parse_result.get("status") != "SUCCESS":
            return {
                "@TYPE": "@ERROR",
                "message": parse_result.get("data", {}).get("message", "자연어 파싱 실패")
            }
        
        parsed_data = parse_result.get("data", {})
        
        # 2단계: 검색 엔진으로 결과 처리
        return self.search_engine.process_query_result(parsed_data)

    # ---------------------------------------------------------
    # [유스케이스 3] 폴더 배치 처리
    # ---------------------------------------------------------
    def process_folder_batch(self, folder_path: str, progress_callback=None) -> Dict[str, Any]:
        """
        폴더 내 파일들을 일괄 처리
        progress_callback: 진행 상황을 전달받을 콜백 함수
        """
        folder_path = self._normalize_path(folder_path)
        
        if not os.path.exists(folder_path):
            return {
                "@TYPE": "@ERROR",
                "message": f"폴더를 찾을 수 없습니다: {folder_path}"
            }

        valid_extensions = (
            '.txt', '.pdf', '.docx', '.xlsx', '.pptx', '.hwp', '.hwpx',
            '.jpg', '.jpeg', '.png', '.webp', '.bmp', '.gif',
            '.mp3', '.mp4', '.wav', '.m4a', '.mkv', '.avi'
        )
        
        files_to_process = []
        
        for root, _, files in os.walk(folder_path):
            for file in files:
                if file.lower().endswith(valid_extensions):
                    full_path = os.path.join(root, file)
                    clean_path = os.path.abspath(os.path.normpath(
                        full_path.replace('￥', '/').replace('\\', '/')))
                    files_to_process.append(clean_path)
        
        if not files_to_process:
            return {
                "@TYPE": "@ERROR",
                "message": "스캔할 지원 파일이 지정된 경로에 없습니다."
            }

        total_count = len(files_to_process)
        success_count = 0
        error_count = 0
        
        # 대량 처리를 위한 DB 세션 시작
        with self.registry.bulk_session():
            for idx, file_path in enumerate(files_to_process, start=1):
                file_name = os.path.basename(file_path)
                
                if progress_callback:
                    progress_callback(f"AI 분석 중 ({idx}/{total_count}): {file_name}")
                
                try:
                    self.process_file_upload(file_path)
                    success_count += 1
                except Exception as e:
                    print(f"파일 처리 실패 ({file_name}): {str(e)}")
                    error_count += 1
        
        # 처리 완료 후 DB 동기화
        removed_files = self.sync_db_with_disk()
        
        return {
            "@TYPE": "@SUCCESS",
            "message": f"폴더 처리 완료: 성공 {success_count}개, 실패 {error_count}개",
            "total_files": total_count,
            "success_count": success_count,
            "error_count": error_count,
            "removed_files": len(removed_files)
        }

    # ---------------------------------------------------------
    # [유틸리티] DB 상태 조회
    # ---------------------------------------------------------
    def get_db_stats(self) -> Dict[str, Any]:
        """DB 현재 상태 통계 정보 반환"""
        conn = self.registry._get_conn()
        try:
            cursor = conn.cursor()
            
            # 전체 파일 수
            cursor.execute("SELECT COUNT(*) FROM files")
            total_files = cursor.fetchone()[0]
            
            # 카테고리별 분포
            cursor.execute("""
                SELECT category, COUNT(*) as count 
                FROM files 
                WHERE category IS NOT NULL AND category != ''
                GROUP BY category
            """)
            category_stats = {row[0]: row[1] for row in cursor.fetchall()}
            
            return {
                "total_files": total_files,
                "category_distribution": category_stats,
                "db_path": self.db_path
            }
        finally:
            if self.registry._bulk_conn is None:
                conn.close()

    # ---------------------------------------------------------
    # [유스케이스 5] 태그 기반 파일 정리
    # ---------------------------------------------------------
    def scan_directory_files(self, directory: str) -> List[Dict[str, Any]]:
        """선택한 디렉터리에서 지원 형식의 파일 목록을 반환합니다."""
        valid_extensions = (
            ".txt", ".pdf", ".docx", ".xlsx", ".pptx", ".hwp", ".hwpx",
            ".jpg", ".jpeg", ".png", ".webp", ".bmp", ".gif",
            ".mp3", ".mp4", ".wav", ".m4a", ".mkv", ".avi",
        )
        directory = self._normalize_path(directory)
        if not os.path.isdir(directory):
            return []

        files: List[Dict[str, Any]] = []
        for root, _, names in os.walk(directory):
            for name in names:
                if name.lower().endswith(valid_extensions):
                    files.append({
                        "file_name": name,
                        "file_path": os.path.join(root, name),
                        "tags": [],
                        "category": "#미분류",
                    })
        return files
    
    def scan_directory_files(self, directory: str) -> List[Dict[str, Any]]:
        """디렉토리에서 지원되는 파일들 스캔"""
        valid_extensions = (
            '.txt', '.pdf', '.docx', '.xlsx', '.pptx', '.hwp', '.hwpx',
            '.jpg', '.jpeg', '.png', '.webp', '.bmp', '.gif',
            '.mp3', '.mp4', '.wav', '.m4a', '.mkv', '.avi'
        )
 
        files = []
        directory = os.path.abspath(directory)
 
        if not os.path.exists(directory):
            return files
 
        for root, _, filenames in os.walk(directory):
            for filename in filenames:
                if filename.lower().endswith(valid_extensions):
                    file_path = os.path.join(root, filename)
                    files.append({
                        "file_name": filename,
                        "file_path": file_path,
                        "tags": [],
                        "category": "#미분류"
                    })
 
        return files
    
    def get_files_for_organize(self) -> List[Dict[str, Any]]:
        """태그가 있는 DB 파일을 정리 화면용 데이터로 조회합니다."""
        conn = self.registry._get_conn()
        try:
            rows = conn.execute(
                """
                SELECT id, file_name, file_path, tags, category
                FROM files
                WHERE tags IS NOT NULL AND tags != ''
                ORDER BY category, file_name
                """
            ).fetchall()
            return [
                {
                    "id": row[0],
                    "file_name": row[1],
                    "file_path": row[2],
                    "tags": row[3].split(",") if row[3] else [],
                    "category": row[4],
                }
                for row in rows
            ]
        finally:
            if self.registry._bulk_conn is None:
                conn.close()

    @staticmethod
    def group_files_by_tags(files: List[Dict[str, Any]]) -> Dict[str, List[Dict[str, Any]]]:
        """첫 번째 태그를 기준으로 파일을 그룹화합니다."""
        groups: Dict[str, List[Dict[str, Any]]] = {}
        for file_info in files:
            tags = file_info.get("tags", [])
            if not tags:
                continue
            tag_name = tags[0].strip().lstrip("#").strip()
            if tag_name:
                groups.setdefault(tag_name, []).append(file_info)
        return groups

    @staticmethod
    def _available_destination(directory: str, file_name: str) -> str:
        """기존 파일을 덮어쓰지 않는 이동 대상 경로를 만듭니다."""
        candidate = os.path.join(directory, file_name)
        if not os.path.exists(candidate):
            return candidate
        stem, extension = os.path.splitext(file_name)
        index = 1
        while os.path.exists(candidate):
            candidate = os.path.join(directory, f"{stem} ({index}){extension}")
            index += 1
        return candidate

    def organize_files(
        self, groups: Dict[str, List[Dict[str, Any]]], base_path: str
    ) -> Dict[str, Any]:
        """태그별 폴더로 파일을 이동하고 DB 경로를 함께 갱신합니다."""
        base_path = self._normalize_path(base_path)
        if not os.path.isdir(base_path):
            return {"success": False, "message": f"기본 경로가 존재하지 않습니다: {base_path}", "errors": []}

        conn = self.registry._get_conn()
        owns_conn = self.registry._bulk_conn is None
        moved_files: List[Dict[str, str]] = []
        errors: List[str] = []
        try:
            conn.execute("BEGIN IMMEDIATE")
            for tag_name, files in groups.items():
                safe_tag = "".join(char for char in tag_name if char not in r'\\/:*?\"<>|').strip()
                if not safe_tag:
                    errors.append(f"사용할 수 없는 태그 이름: {tag_name}")
                    continue
                target_dir = os.path.join(base_path, safe_tag)
                os.makedirs(target_dir, exist_ok=True)
                for file_info in files:
                    old_path = self._normalize_path(file_info.get("file_path", ""))
                    if not os.path.isfile(old_path):
                        errors.append(f"파일 없음: {old_path}")
                        continue
                    new_path = self._available_destination(target_dir, os.path.basename(old_path))
                    try:
                        shutil.move(old_path, new_path)
                        conn.execute(
                            "UPDATE files SET file_name = ?, file_path = ?, source_path = ? WHERE id = ?",
                            (os.path.basename(new_path), new_path, target_dir, file_info["id"]),
                        )
                        moved_files.append({"old_path": old_path, "new_path": new_path, "tag": safe_tag})
                    except OSError as exc:
                        errors.append(f"이동 실패 ({old_path}): {exc}")
            conn.commit()
            return {
                "success": True,
                "message": f"파일 정리 완료: {len(moved_files)}개 파일 이동, {len(errors)}개 오류",
                "moved_files": moved_files,
                "errors": errors,
            }
        except Exception as exc:
            conn.rollback()
            return {"success": False, "message": f"파일 정리 실패: {exc}", "errors": [str(exc)]}
        finally:
            if owns_conn:
                conn.close()

# =========================================================
# 하위 호환성을 위한 별칭 클래스
# =========================================================
class MainProcessor(ClasqCore):
    """기존 코드와의 호환성을 위한 별칭 클래스"""
    pass
#========================core.py===========================

#========================db_manager.py=====================
# =========================================================
# [db_manager.py]
# DB 자동 저장 무결성 보장, 파일 해시 기반 중복 감지,
# 물리적 파일 동기화 정책을 전담하는 모듈 (담당: 팀원 C)
#
# 기존 main_processor.py에 있던 _init_db / _save_to_db 로직을
# 이 클래스로 완전히 이전하고, 아래 3가지를 보강했다.
#   1) 파일 해시(SHA-256) 기반 중복 감지 -> 내용이 같은 파일은
#      file_path가 달라도 걸러내고 정책에 따라 물리적으로 격리
#   2) 대량 스캔(폴더 전체) 시 커넥션 재사용 + WAL 모드로
#      "database is locked" 및 성능 저하 방지
#   3) DB 레코드와 실제 디스크 파일 상태를 맞추는 동기화 기능
#      (외부에서 파일이 지워지거나 이동된 경우 정리)
# =========================================================
import os
import sqlite3
import hashlib
import shutil
import time
from contextlib import contextmanager
from typing import Dict, Any, Optional, List
import re


class FileRegistryManager:
    """
    [DB / 물리 파일 담당 클래스]

    main_processor.py, gui_app.py 양쪽에서 공통으로 사용하는
    '파일 레코드 저장소'. files 테이블 스키마와 저장/중복감지/동기화
    로직을 이 클래스 하나로 단일화하여, main_processor.py와
    gui_app.py에 스키마가 중복 정의되던 문제를 제거한다.
    """

    def __init__(
        self,
        db_path: str = "file_manager.db",
        duplicate_policy: str = "quarantine",
        duplicates_dir_name: str = "_duplicates",
    ):
        """
        duplicate_policy (내용이 동일한 파일을 새로 발견했을 때의 처리 정책)
          - "quarantine" : 스캔된 파일을 같은 폴더의 _duplicates 하위 폴더로
                           물리적으로 옮기고, DB에는 원본 파일 경로만 유지 (기본값)
          - "skip"       : 물리 파일은 그대로 두고, DB에도 새로 기록하지 않음
          - "keep"       : 중복이어도 그냥 완전히 별개 파일처럼 DB에 저장
        """
        self.db_path = db_path
        self.duplicate_policy = duplicate_policy
        self.duplicates_dir_name = duplicates_dir_name
        self._bulk_conn: Optional[sqlite3.Connection] = None  # 대량 스캔용 재사용 커넥션
        self._init_db()

    # ---------------------------------------------------------
    # 0. 커넥션 helper
    #    WAL 모드: GUI 스레드가 읽는 동안 워커 스레드가 써도 서로 안 막힘
    #    busy_timeout: 스캔 중 순간적인 락 충돌 시 즉시 에러 대신 잠깐 대기
    # ---------------------------------------------------------
    def _connect(self) -> sqlite3.Connection:
        conn = sqlite3.connect(self.db_path, timeout=30)
        conn.execute("PRAGMA journal_mode=WAL")
        conn.execute("PRAGMA busy_timeout=30000")
        return conn

    def _get_conn(self) -> sqlite3.Connection:
        """bulk_session() 안이면 재사용 커넥션을, 아니면 새 커넥션을 반환"""
        return self._bulk_conn if self._bulk_conn is not None else self._connect()

    # ---------------------------------------------------------
    # 1. DB 초기화 및 스키마 마이그레이션
    #    (기존에 이미 만들어진 files.db가 있어도 안전하게 컬럼만 추가)
    # ---------------------------------------------------------
    def _init_db(self) -> None:
        conn = self._connect()
        try:
            conn.execute(
                """
                CREATE TABLE IF NOT EXISTS files (
                    id INTEGER PRIMARY KEY AUTOINCREMENT,
                    file_name TEXT,
                    file_path TEXT UNIQUE,
                    ai_comment TEXT,
                    category TEXT
                )
                """
            )
            self._migrate_add_columns(
                conn,
                {
                    "file_hash": "TEXT",
                    "file_size": "INTEGER",
                    "created_at": "TEXT",
                    "updated_at": "TEXT",
                    "tags": "TEXT",
                    "source_path": "TEXT",
                },
            )
            conn.execute(
                "CREATE INDEX IF NOT EXISTS idx_files_hash ON files(file_hash)")
            
            # 경로 관리 테이블 생성
            conn.execute(
                """
                CREATE TABLE IF NOT EXISTS managed_paths (
                    id INTEGER PRIMARY KEY AUTOINCREMENT,
                    path TEXT UNIQUE,
                    created_at TEXT DEFAULT CURRENT_TIMESTAMP
                )
                """
            )
            conn.commit()
        finally:
            conn.close()

    def _migrate_add_columns(self, conn: sqlite3.Connection, columns: Dict[str, str]) -> None:
        """기존 DB에 신규 컬럼이 없으면 ALTER TABLE로 안전하게 추가 (기존 데이터는 그대로 보존)"""
        existing_cols = {row[1] for row in conn.execute(
            "PRAGMA table_info(files)").fetchall()}
        for col_name, col_type in columns.items():
            if col_name not in existing_cols:
                conn.execute(
                    f"ALTER TABLE files ADD COLUMN {col_name} {col_type}")

    # ---------------------------------------------------------
    # 2. 대량 스캔용 커넥션 재사용 세션
    #    - 파일마다 connect/close를 반복하지 않아 폴더 스캔 성능이 개선됨
    #    - 각 파일 저장은 여전히 개별 트랜잭션(BEGIN~COMMIT)으로 처리되므로
    #      중간에 한 파일이 실패해도 그 전까지 저장된 파일들은 유지된다
    #      (전체를 한 트랜잭션으로 묶지 않는 이유: 수백 개 중 1개 실패로
    #       이미 분석 끝난 나머지 파일들의 DB 기록까지 날아가면 안 되기 때문)
    # ---------------------------------------------------------
    @contextmanager
    def bulk_session(self):
        conn = self._connect()
        self._bulk_conn = conn
        try:
            yield self
        finally:
            conn.close()
            self._bulk_conn = None

    # ---------------------------------------------------------
    # 3. 파일 해시 계산 (중복 감지의 기준)
    # ---------------------------------------------------------
    @staticmethod
    def compute_file_hash(file_path: str, chunk_size: int = 65536) -> str:
        sha256 = hashlib.sha256()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(chunk_size), b""):
                sha256.update(chunk)
        return sha256.hexdigest()

    def _find_by_hash(
        self, conn: sqlite3.Connection, file_hash: str, exclude_path: Optional[str] = None
    ) -> Optional[tuple]:
        query = "SELECT file_path FROM files WHERE file_hash = ?"
        params: List[Any] = [file_hash]
        if exclude_path:
            query += " AND file_path != ?"
            params.append(exclude_path)
        return conn.execute(query, params).fetchone()

    # ---------------------------------------------------------
    # 4. 물리적 파일 동기화 정책
    # ---------------------------------------------------------
    def _quarantine_file(self, file_path: str) -> str:
        """내용이 같은 파일을 같은 폴더의 _duplicates 하위 폴더로 이동시키고 새 경로 반환"""
        base_dir = os.path.dirname(file_path) or "."
        dest_dir = os.path.join(base_dir, self.duplicates_dir_name)
        os.makedirs(dest_dir, exist_ok=True)

        filename = os.path.basename(file_path)
        dest_path = os.path.join(dest_dir, filename)

        stem, ext = os.path.splitext(filename)
        counter = 1
        while os.path.exists(dest_path):
            dest_path = os.path.join(dest_dir, f"{stem}({counter}){ext}")
            counter += 1

        shutil.move(file_path, dest_path)
        return dest_path

    def sync_missing_files(self) -> List[str]:
        """
        DB에는 있지만 실제 디스크에는 없는(사용자가 탐색기에서 지우거나 옮긴)
        레코드를 찾아 정리한다. '물리적 파일 동기화'의 핵심 기능.
        반환값: 정리(삭제)된 file_path 리스트
        """
        conn = self._get_conn()
        owns_conn = self._bulk_conn is None
        removed: List[str] = []
        try:
            rows = conn.execute("SELECT file_path FROM files").fetchall()
            for (file_path,) in rows:
                if not os.path.exists(file_path):
                    conn.execute(
                        "DELETE FROM files WHERE file_path = ?", (file_path,))
                    removed.append(file_path)
            conn.commit()
        finally:
            if owns_conn:
                conn.close()
        return removed

    # ---------------------------------------------------------
    # 5. 경로 관리 함수
    # ---------------------------------------------------------
    def add_managed_path(self, path: str) -> Dict[str, Any]:
        """관리할 경로를 DB에 추가"""
        path = os.path.abspath(path)
        if not os.path.exists(path):
            return {"success": False, "message": f"경로가 존재하지 않습니다: {path}"}
        
        conn = self._get_conn()
        owns_conn = self._bulk_conn is None
        try:
            conn.execute(
                "INSERT OR IGNORE INTO managed_paths (path) VALUES (?)",
                (path,)
            )
            conn.commit()
            return {"success": True, "message": f"경로가 추가되었습니다: {path}"}
        except Exception as e:
            return {"success": False, "message": f"경로 추가 실패: {str(e)}"}
        finally:
            if owns_conn:
                conn.close()

    def get_managed_paths(self) -> List[str]:
        """관리 중인 모든 경로 조회"""
        conn = self._get_conn()
        owns_conn = self._bulk_conn is None
        try:
            rows = conn.execute("SELECT path FROM managed_paths ORDER BY created_at").fetchall()
            return [row[0] for row in rows]
        finally:
            if owns_conn:
                conn.close()

    def remove_managed_path(self, path: str) -> Dict[str, Any]:
        """관리 경로 제거"""
        path = os.path.abspath(path)
        conn = self._get_conn()
        owns_conn = self._bulk_conn is None
        try:
            cursor = conn.execute(
                "DELETE FROM managed_paths WHERE path = ?",
                (path,)
            )
            conn.commit()
            if cursor.rowcount > 0:
                return {"success": True, "message": f"경로가 제거되었습니다: {path}"}
            else:
                return {"success": False, "message": f"경로를 찾을 수 없습니다: {path}"}
        except Exception as e:
            return {"success": False, "message": f"경로 제거 실패: {str(e)}"}
        finally:
            if owns_conn:
                conn.close()

    # ---------------------------------------------------------
    # 6. 메인 저장 함수 (main_processor.py의 _save_to_db 대체)
    # ---------------------------------------------------------
    def save_file_result(self, file_path: str, metadata_result: Dict[str, Any]) -> Dict[str, Any]:
        """
        AI 분석 완료 후 호출되는 저장 함수.
          1) 파일 해시 계산
          2) 동일 해시(내용 동일) 레코드가 이미 있는지 확인 -> 중복 정책 적용
          3) UPSERT(신규 INSERT / 기존 UPDATE)로 DB 반영
        한 번의 호출은 항상 하나의 트랜잭션으로 처리되어, 저장 도중 오류가 나도
        DB에 절반만 반영된 깨진 레코드가 남지 않는다.
        """
        result: Dict[str, Any] = {"success": False,
                                  "file_path": file_path, "is_duplicate": False}

        if not os.path.exists(file_path):
            result["message"] = f"파일을 찾을 수 없습니다: {file_path}"
            return result

        conn = self._get_conn()
        owns_conn = self._bulk_conn is None

        try:
            conn.execute("BEGIN IMMEDIATE")

            file_hash = self.compute_file_hash(file_path)
            dup_row = self._find_by_hash(
                conn, file_hash, exclude_path=file_path)

            final_path = file_path
            if dup_row is not None:
                result["is_duplicate"] = True
                result["duplicate_of"] = dup_row[0]

                if self.duplicate_policy == "quarantine":
                    final_path = self._quarantine_file(file_path)
                    result["quarantined_to"] = final_path
                elif self.duplicate_policy == "skip":
                    conn.commit()
                    result["success"] = True
                    result["message"] = "중복 파일 - DB 반영 없이 스킵됨"
                    return result
                # "keep" 정책이면 별개 파일처럼 그대로 진행

            meta = metadata_result.get("metadata", {})
            file_name = os.path.basename(final_path)
            ai_comment = meta.get("ai_comment", "")
            tags = meta.get("tags", [])
            category = f"#{tags[0]}" if tags else "#일반"
            tags_str = ",".join(tags) if tags else ""
            source_path = os.path.dirname(final_path)
            file_size = os.path.getsize(final_path)
            now = time.strftime("%Y-%m-%d %H:%M:%S")

            # 🌟 [개선] INSERT OR REPLACE -> UPSERT(ON CONFLICT)로 변경
            #    기존 INSERT OR REPLACE는 내부적으로 DELETE 후 INSERT를 수행해
            #    기존 레코드의 id가 매번 바뀌는 부작용이 있었다. ON CONFLICT
            #    DO UPDATE는 같은 id를 유지한 채 값만 갱신하므로 향후 다른
            #    테이블에서 file_id를 참조(FK)하게 되어도 안전하다.
            conn.execute(
                """
                INSERT INTO files (file_name, file_path, ai_comment, category,
                                    file_hash, file_size, created_at, updated_at, tags, source_path)
                VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
                ON CONFLICT(file_path) DO UPDATE SET
                    file_name = excluded.file_name,
                    ai_comment = excluded.ai_comment,
                    category = excluded.category,
                    file_hash = excluded.file_hash,
                    file_size = excluded.file_size,
                    updated_at = excluded.updated_at,
                    tags = excluded.tags,
                    source_path = excluded.source_path
                """,
                (file_name, final_path, ai_comment,
                 category, file_hash, file_size, now, now, tags_str, source_path),
            )

            conn.commit()
            result["success"] = True
            result["file_path"] = final_path
            return result

        except Exception as e:
            conn.rollback()
            result["message"] = f"DB 저장 오류: {e}"
            return result
        finally:
            if owns_conn:
                conn.close()

    # ---------------------------------------------------------
    # 6. DB 초기화 (gui_app.py의 reset_db_and_path 대체용)
    #    - os.remove() 대신 SQL DELETE + sqlite_sequence 리셋 방식을
    #      그대로 유지 (WinError 32 프로세스 점유 에러 방지 정책 계승)
    #    - 스키마 정의를 이 클래스 한 곳에서만 관리하도록 일원화
    # ---------------------------------------------------------
    def reset_all(self) -> None:
        conn = self._get_conn()
        owns_conn = self._bulk_conn is None
        try:
            conn.execute("DELETE FROM files;")
            conn.execute("DELETE FROM sqlite_sequence WHERE name='files';")
            conn.commit()
        finally:
            if owns_conn:
                conn.close()

    def organize_file(self, file_id: int, target_category: str, base_dir: str) -> bool:
        """물리적 파일 이동 및 DB 업데이트"""
        # 0. 폴더명으로 쓸 수 없는 문자 제거
        safe_category = re.sub(r'[\\/*?:"<>|]', "", target_category).strip()

        conn = self._get_conn()
        try:
            row = conn.execute(
                "SELECT file_path FROM files WHERE id = ?", (file_id,)
            ).fetchone()

            if not row:
                return False

            old_path = row[0]

            if not os.path.exists(old_path):
                print(f"[에러] 원본 파일을 찾을 수 없음: {old_path}")
                return False

            new_dir = os.path.join(base_dir, safe_category)
            os.makedirs(new_dir, exist_ok=True)
            new_path = os.path.join(new_dir, os.path.basename(old_path))

            # 파일 이동
            shutil.move(old_path, new_path)

            # DB 업데이트
            conn.execute(
                "UPDATE files SET file_path = ? WHERE id = ?", (
                    new_path, file_id)
            )
            conn.commit()
            return True

        except Exception as e:
            print(f"[파일 이동 실패]: {e}")
            return False

        finally:
            if self._bulk_conn is None:  # bulk_session이 아닐 때만 닫음
                conn.close()

    def delete_file(self, file_id: int) -> bool:
        """
        [물리적 파일 삭제 및 DB 레코드 제거]
        지정된 file_id의 실제 파일을 디스크 및 DB에서 삭제.
        """
        conn = self._get_conn()
        try:
            # 1. DB에서 파일 경로 찾기
            row = conn.execute(
                "SELECT file_path FROM files WHERE id = ?", (file_id,)).fetchone()
            if not row:
                return False

            target_path = row[0]

            # 2. 물리적 디스크에서 파일 삭제 (파일이 존재할 경우만)
            if os.path.exists(target_path):
                os.remove(target_path)

            # 3. DB 레코드 삭제
            conn.execute("DELETE FROM files WHERE id = ?", (file_id,))
            conn.commit()
            return True

        except Exception as e:
            print(f"[파일 삭제 실패]: {e}")
            return False
        finally:
            if self._bulk_conn is None:
                conn.close()

    def rename_file(self, file_id: int, new_filename: str) -> bool:
        """
        [물리적 파일 이름 변경 및 DB 갱신]
        지정된 file_id의 실제 파일 이름을 바꾸고, DB의 file_name과 file_path를 업데이트.
        (new_filename에는 확장자도 포함되어야 함. 예: "휴가계획.txt")
        """
        # 윈도우에서 사용할 수 없는 특수문자 방어
        safe_name = re.sub(r'[\\/*?:"<>|]', "", new_filename).strip()

        conn = self._get_conn()
        try:
            # 1. DB에서 현재 파일 경로 찾기
            row = conn.execute(
                "SELECT file_path FROM files WHERE id = ?", (file_id,)).fetchone()
            if not row:
                return False

            old_path = row[0]
            if not os.path.exists(old_path):
                print(f"[에러] 원본 파일을 찾을 수 없음: {old_path}")
                return False

            # 2. 새로운 경로(이름) 생성
            base_dir = os.path.dirname(old_path)
            new_path = os.path.join(base_dir, safe_name)

            # 3. 물리적 파일 이름 변경
            os.rename(old_path, new_path)

            # 4. DB 정보 업데이트 (파일 이름, 파일 경로 둘 다 변경)
            conn.execute(
                "UPDATE files SET file_name = ?, file_path = ? WHERE id = ?",
                (safe_name, new_path, file_id)
            )
            conn.commit()
            return True

        except Exception as e:
            print(f"[파일 이름 변경 실패]: {e}")
            return False
        finally:
            if self._bulk_conn is None:
                conn.close()
#========================db_manager.py=====================
#========================file_pipeline.py=====================
# =========================================================
# [file_pipeline.py] 
# 문서, 이미지, 미디어 파일 데이터 추출/전처리 및 
# 로컬 AI(Ollama) 메타데이터 생성 모듈
# =========================================================

# =====================================================================
# [예외 처리 및 방어적 프로그래밍 아키텍처 요약]
# 1. 사전 유효성 및 4대 방어막 검사
#    - 파일 존재 여부 및 0 바이트 빈 파일 예외 차단
#    - 허용되지 않은 압축 파일(ZIP, RAR 등) 진입 원천 거부
#    - 대용량 파일 용량 제한(MB 단위) 초과 검사 및 구버전 오피스 파일(.doc, .xls, .ppt) 지원 제외
# 2. 포맷별 맞춤형 예외 방어 (TextExtractor)
#    - 이미지: 애니메이션 GIF 다중 프레임 안전 처리, 이중 verify() 버그 해결, 해상도 폭탄(DecompressionBomb) 방어
#    - 텍스트/문서: UTF-8 디코딩 실패 시 CP949(EUC-KR) 자동 폴백, 암호화된 PDF 및 손상된 HWPX 압축 파일 예외 처리
# 3. AI 통신 및 파싱 방어 (FileAnalyzer)
#    - Ollama API 통신 중 네트워크 단절(ConnectionError), 타임아웃(Timeout), JSON 파싱 깨짐(JSONDecodeError) 세분화 예외 처리
#    - AI 분석 실패 시 프로그램 종료 방지 및 최소 메타데이터를 담은 폴백(Fallback) 응답(JSON) 반환
# 4. 전체 프로세스 중단 방지 (Graceful Degradation)
#    - 최상단 extract() 메서드에 거대한 try-except 블록을 두어 파싱 오류 시에도 파일명 기반 대체 텍스트로 태깅 프로세스 완주 보장
# =====================================================================

# 1. 파이썬 표준 라이브러리 (기능별)
import os          # 파일 경로 확인, 파일 크기 검사, 확장자 추출 모듈
import re          # 정규표현식(특수문자 제거, 패턴 매칭) 모듈
import json        # JSON 문자열 변환 및 데이터 파싱 모듈
import io          # 메모리 내 바이너리 바이트 버퍼 처리 모듈
import zlib        # HWP 파일 데이터 압축 해제(Decompress) 모듈
import zipfile     # HWPX/DOCX 등 ZIP 포맷 압축 해제 모듈
import xml.etree.ElementTree as ET  # XML 구조 파일 텍스트 추출용 모듈
import requests    # 예외 타입 처리용 라이브러리
from ollama_manager import OllamaManager
from datetime import datetime       # 타임스탬프(분석 시간) 기록용 모듈
from typing import Dict, Any, Tuple # 파이썬 함수 리턴 타입 명시용 모듈

# 2. 문서 및 이미지 파싱용 외부 제3자 라이브러리
from pypdf import PdfReader                  # PDF 파일 텍스트 추출
from docx import Document                    # Word(.docx) 파일 문단 텍스트 추출
import openpyxl                              # Excel(.xlsx) 셀 데이터 추출
from PIL import Image, UnidentifiedImageError # 이미지 리사이징 및 손상 검사

# 3. 외부 선택 설치 패키지 안전검사 (try-except 모듈 동적 로딩)
# olefile (구버전 HWP 바이너리 파싱용)
try:
    import olefile
    HAS_OLEFILE = True
except ImportError:
    HAS_OLEFILE = False

# python-pptx (PowerPoint 파싱용)
try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# openai-whisper (음성 인식 STT 전용)
try:
    import whisper
    HAS_WHISPER = True
except ImportError:
    HAS_WHISPER = False


# =========================================================
# [커스텀 예외 클래스]
# =========================================================
class FilePreprocessError(Exception):
    """파일 전처리(읽기/해석) 과정에서 오류가 발생했을 때 던지는 예외 클래스"""
    pass


# =========================================================
# [Step 1] 파일 원문/데이터 추출 클래스 (TextExtractor)
# =========================================================
class TextExtractor:
    """확장자별 문서, 이미지, 오디오/비디오 데이터 추출 및 6중 방어막 예외 처리 클래스"""

    # A. 처리 가능한 파일 확장자 그룹 정의
    IMAGE_EXTENSIONS = ('.jpg', '.jpeg', '.png', '.webp', '.bmp', '.gif', '.tiff', '.tif')
    
    DOC_EXTENSIONS = (
        '.txt', '.pdf', '.docx', '.xlsx', '.pptx', 
        '.hwp', '.hwpx',
        '.csv', '.json', '.xml', '.yaml', '.yml', 
        '.html', '.htm', '.md', '.markdown'
    )
    
    AUDIO_VIDEO_EXTENSIONS = ('.mp3', '.wav', '.m4a', '.mp4', '.mkv', '.avi')

    # B. 보안 및 오류 방지를 위해 명시적으로 시스템 입력을 차단할 압축 파일
    ARCHIVE_EXTENSIONS = ('.zip', '.rar', '.7z', '.tar', '.gz', '.bz2', '.iso')

    def __init__(
        self, 
        max_chars: int = 2000, 
        max_img_size: int = 512, 
        whisper_model_name: str = "base",
        max_doc_size_mb: int = 50,       # 문서 파일 최대 50MB 용량 제한
        max_img_size_mb: int = 100,      # 이미지 파일 최대 100MB 용량 제한
        max_media_size_mb: int = 500     # 오디오/비디오 최대 500MB 용량 제한
    ):
        """기본 제한 설정(글자 수 제한, 최대 용량 제한 등)을 초기화하는 생성자"""
        self.max_chars = max_chars
        self.max_img_size = max_img_size
        self.whisper_model_name = whisper_model_name
        
        # MB 바이트 단위로 계산 변환 (1MB = 1024 * 1024 Bytes)
        self.max_doc_bytes = max_doc_size_mb * 1024 * 1024
        self.max_img_bytes = max_img_size_mb * 1024 * 1024
        self.max_media_bytes = max_media_size_mb * 1024 * 1024
        
        self._whisper_model = None  # Whisper 모델 지연 로딩(필요할 때 메모리 적재)

    def _sanitize_text(self, text: str) -> str:
        """문서 내부의 깨진 문자나 시스템 특수 제어 문자를 안전하게 제거하는 정제 함수"""
        if not text:
            return ""
        return re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', text)

    def is_image_file(self, file_path: str) -> bool:
        """해당 파일 경로가 지원되는 이미지 확장자인지 여부 확인"""
        ext = os.path.splitext(file_path)[1].lower()
        return ext in self.IMAGE_EXTENSIONS

    def is_media_file(self, file_path: str) -> bool:
        """해당 파일 경로가 지원되는 오디오/비디오 확장자인지 여부 확인"""
        ext = os.path.splitext(file_path)[1].lower()
        return ext in self.AUDIO_VIDEO_EXTENSIONS

    # ---------------------------------------------------------
    # [1] 이미지 파일 전처리 함수 (애니메이션 GIF 및 확장자 불일치 예외 완벽 방어)
    # ---------------------------------------------------------
    def process_image(self, file_path: str) -> Tuple[bytes, str]:
        """이미지 파일 유효성 검사, 썸네일 축소 및 메모리 바이너리(Bytes) 변환 함수"""
        
        # 1-1. 존재 여부 검사
        if not os.path.exists(file_path):
            return b"", "ERROR: 존재하지 않는 이미지 파일입니다."

        # 1-2. 0 바이트 빈 파일 검사
        file_size = os.path.getsize(file_path)
        if file_size == 0:
            return b"", "ERROR: 내용이 없는 빈 이미지 파일(0 Byte)입니다."

        # 1-3. 용량 제한 초과 검사
        if file_size > self.max_img_bytes:
            max_mb = self.max_img_bytes // (1024 * 1024)
            return b"", f"ERROR: 이미지 용량이 제한({max_mb}MB)을 초과했습니다. ({file_size / (1024*1024):.1f}MB)"

        try:
            # 💡 [버그 수정] img.verify() 이중 호출 제거: verify()를 부르면 파일 포인터가 손상되어
            # 특정 WEBP, GIF, 확장자가 변형된 이미지 파일에서 UnidentifiedImageError가 터지던 문제를 완벽 해결했습니다.
            with Image.open(file_path) as img:
                # 애니메이션 GIF 등의 다중 프레임 이미지일 경우 첫 번째 대표 프레임 선택
                if getattr(img, "is_animated", False):
                    img.seek(0)

                img_rgb = img.convert("RGB") # Vision AI 전달용 RGB 포맷 통합
                img_rgb.thumbnail((self.max_img_size, self.max_img_size)) # 축소 리사이징

                buffer = io.BytesIO()
                img_rgb.save(buffer, format="JPEG", quality=80) # JPEG 압축 저장
                return buffer.getvalue(), "SUCCESS"

        except PermissionError:
            return b"", "ERROR: 이미지 파일 접근 권한이 없거나 다른 프로그램에서 사용 중입니다."
        except Image.DecompressionBombError:
            return b"", "ERROR: 해상도가 너무 큰 비정상적인 이미지(DecompressionBomb)입니다."
        except Exception as e:
            # 가짜 확장자나 헤더 손상 시에도 전체 스캔이 중단되지 않도록 예외 차단 후 메시지 반환
            return b"", f"ERROR: 이미지 읽기 및 변환 실패 ({str(e)})"

    # ---------------------------------------------------------
    # [2] 오디오/비디오 미디어 파일 전처리 (FFmpeg 미설치/STT 오류 안전 방어)
    # ---------------------------------------------------------
    def process_media(self, file_path: str) -> Tuple[str, str]:
        """Whisper AI 모델을 활용해 음성을 텍스트로 변환하는 함수"""
        file_name = os.path.basename(file_path)
        
        # 2-1. 패키지 설치 여부 검사
        if not HAS_WHISPER:
            return f"미디어 파일: {file_name} (openai-whisper 미설치)", "SUCCESS"

        # 2-2. 파일 존재 여부 및 용량 검사
        if not os.path.exists(file_path) or os.path.getsize(file_path) == 0:
            return f"미디어 파일: {file_name} (빈 파일 또는 존재하지 않음)", "SUCCESS"

        try:
            # Whisper 모델 지연 로딩
            if self._whisper_model is None:
                self._whisper_model = whisper.load_model(self.whisper_model_name)

            # STT 음성 추출 실행
            result = self._whisper_model.transcribe(file_path)
            extracted_text = result.get("text", "").strip()

            if not extracted_text:
                extracted_text = f"미디어 파일명: {file_name} (음성 인식 데이터 없음)"

            clean_text = self._sanitize_text(extracted_text)
            return clean_text[:self.max_chars].strip(), "SUCCESS"

        except Exception as e:
            # FFmpeg가 미설치되어 있거나 디코딩이 실패해도 파일명을 텍스트로 넘겨 AI 분석 진행!
            return f"음성 추출 불가 미디어 파일: {file_name} (STT 실패: {str(e)})", "SUCCESS"

    # ---------------------------------------------------------
    # [3] 일반 문서 및 데이터 파일 텍스트 추출 메인 함수 (파싱 오류 완벽 방어)
    # ---------------------------------------------------------
    def extract(self, file_path: str) -> Tuple[str, str]:
        """확장자별 개별 추출 알고리즘을 호출하고 예외 상황을 일괄 제어하는 메인 함수"""
        
        if not os.path.exists(file_path):
            return "", "ERROR: 존재하지 않는 파일입니다."

        file_size = os.path.getsize(file_path)
        if file_size == 0:
            return "", "ERROR: 내용이 없는 빈 파일(0 Byte)입니다."

        ext = os.path.splitext(file_path)[1].lower()

        # 방어막 1: 단순 압축 파일 명시적 거부
        if ext in self.ARCHIVE_EXTENSIONS:
            return "", f"ERROR: 압축 파일({ext})은 지원하지 않습니다. 압축을 해제한 후 개별 파일로 업로드해 주세요."

        # 방어막 2: 문서 용량 제한 검사
        if file_size > self.max_doc_bytes:
            max_mb = self.max_doc_bytes // (1024 * 1024)
            return "", f"ERROR: 문서 용량이 제한({max_mb}MB)을 초과했습니다. ({file_size / (1024*1024):.1f}MB)"

        # 방어막 3: 구버전 MS 오피스 파일(.doc, .xls, .ppt) 지원 제외
        if ext in ['.doc', '.xls', '.ppt']:
            return "", f"ERROR: 구버전 오피스 파일({ext})은 지원하지 않습니다. 최신 포맷(.docx, .xlsx, .pptx)으로 변환해 주세요."

        try:
            # 포맷별 개별 파서 분기 연결
            if ext in ['.txt', '.csv', '.json', '.xml', '.yaml', '.yml', '.html', '.htm', '.md', '.markdown']:
                text = self._read_txt(file_path)
            elif ext == '.pdf':
                text = self._read_pdf(file_path)
            elif ext == '.docx':
                text = self._read_docx(file_path)
            elif ext == '.xlsx':
                text = self._read_xlsx(file_path)
            elif ext == '.pptx':
                text = self._read_pptx(file_path)
            elif ext == '.hwpx':
                text = self._read_hwpx(file_path)
            elif ext == '.hwp':
                text = self._read_hwp(file_path)
            else:
                return "", f"ERROR: 지원하지 않는 파일 확장자입니다 ({ext})"

            clean_text = self._sanitize_text(text)

            # 💡 텍스트 추출이 빈 값이어도 파일명 자체를 전달하여 AI 태깅 수행
            if not clean_text or not clean_text.strip():
                clean_text = f"문서 파일명: {os.path.basename(file_path)} (내부 텍스트 내용 없음)"

            return clean_text[:self.max_chars].strip(), "SUCCESS"

        except Exception as e:
            # 💡 [버그 수정] HWP/PDF/DOCX 등 파싱 중 FilePreprocessError 등의 예외가 터져도
            # 백엔드 스레드가 멈추지 않고 파일명으로 대체하여 AI 태깅이 수월하게 진행되도록 최상단에서 방어합니다.
            file_name = os.path.basename(file_path)
            return f"문서 파일명: {file_name} (내부 데이터 해석 실패: {str(e)})", "SUCCESS"

    # --- 포맷별 하위 파싱 메서드 모음 ---

    def _read_txt(self, path: str) -> str:
        """대용량 TXT 방어: 최대 1MB 스트리밍 읽기 (UTF-8 / CP949 자동 적용)"""
        try:
            with open(path, 'r', encoding='utf-8') as f:
                return f.read(1024 * 1024)
        except UnicodeDecodeError:
            try:
                with open(path, 'r', encoding='cp949', errors='ignore') as f:
                    return f.read(1024 * 1024)
            except Exception as e:
                raise FilePreprocessError(f"텍스트 인코딩 읽기 실패: {str(e)}")

    def _read_pdf(self, path: str) -> str:
        """PDF 문서의 각 페이지에서 텍스트 읽기 및 암호화 여부 체크"""
        try:
            reader = PdfReader(path)
            if reader.is_encrypted:
                raise FilePreprocessError("암호로 보호된 PDF 파일입니다.")

            extracted_text = []
            total_len = 0

            for page in reader.pages:
                page_text = page.extract_text() or ""
                if page_text.strip():
                    extracted_text.append(page_text)
                    total_len += len(page_text)

                if total_len >= self.max_chars:
                    break

            return "\n".join(extracted_text)
        except FilePreprocessError:
            raise
        except Exception as e:
            raise FilePreprocessError(f"PDF 파싱 실패 ({str(e)})")

    def _read_docx(self, path: str) -> str:
        """Word 문서(.docx)의 각 문단 Paragraph 텍스트 추출"""
        try:
            doc = Document(path)
            full_text = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n".join(full_text)
        except Exception as e:
            raise FilePreprocessError(f"Word 파일 파싱 실패 ({str(e)})")

    def _read_xlsx(self, path: str) -> str:
        """Excel(.xlsx) 메모리 절약 모드(read_only)로 시트 내 셀 데이터 읽기"""
        wb = None
        try:
            wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
            text_list = []
            current_len = 0

            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_str = " ".join([str(cell) for cell in row if cell is not None])
                    if row_str.strip():
                        text_list.append(row_str)
                        current_len += len(row_str)

                    if current_len >= self.max_chars:
                        break
                if current_len >= self.max_chars:
                    break

            return "\n".join(text_list)
        except Exception as e:
            raise FilePreprocessError(f"Excel 파일 파싱 실패 ({str(e)})")
        finally:
            if wb:
                wb.close()

    def _read_pptx(self, path: str) -> str:
        """PowerPoint(.pptx) 슬라이드 도형 내 텍스트 파싱"""
        if not HAS_PPTX:
            raise FilePreprocessError("python-pptx 패키지가 설치되지 않았습니다. ('pip install python-pptx' 필요)")

        try:
            prs = Presentation(path)
            text_list = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_list.append(shape.text.strip())
            return "\n".join(text_list)
        except Exception as e:
            raise FilePreprocessError(f"PPTX 파싱 실패 ({str(e)})")

    def _read_hwpx(self, path: str) -> str:
        """HWPX(ZIP 압축 해제) 형태의 Contents/section.xml 본문 추출"""
        try:
            text_list = []
            with zipfile.ZipFile(path, 'r') as zf:
                section_files = sorted([f for f in zf.namelist() if re.match(r'Contents/section\d+\.xml', f)])
                
                for sec_file in section_files:
                    xml_data = zf.read(sec_file)
                    root = ET.fromstring(xml_data)
                    
                    for elem in root.iter():
                        if elem.tag.endswith('t') and elem.text:
                            text_list.append(elem.text.strip())
                            
            return "\n".join(text_list)
        except zipfile.BadZipFile:
            raise FilePreprocessError("올바른 HWPX 압축 파일이 아니거나 손상되었습니다.")
        except Exception as e:
            raise FilePreprocessError(f"HWPX 파싱 실패 ({str(e)})")

    def _read_hwp(self, path: str) -> str:
        """HWP 오피스 바이너리 문서 OLE 구조에서 BodyText 비트 스트림 복호화 및 추출"""
        if not HAS_OLEFILE:
            raise FilePreprocessError("olefile 패키지가 설치되지 않았습니다. ('pip install olefile' 필요)")

        try:
            ole = olefile.OleFileIO(path)
            dirs = ole.listdir()

            body_sections = [d for d in dirs if d[0] == 'BodyText']
            body_sections.sort()

            if not body_sections:
                ole.close()
                raise FilePreprocessError("HWP 문서 내 본문(BodyText) 영역을 찾을 수 없습니다.")

            text_list = []
            for section in body_sections:
                stream = ole.openstream(section).read()
                
                # Deflate 압축 데이터 해제
                try:
                    decompressed = zlib.decompress(stream, -15)
                except zlib.error:
                    decompressed = stream

                # Record 구조 단위 오프셋 탐색
                i = 0
                while i < len(decompressed):
                    if i + 4 > len(decompressed):
                        break
                    record_header = int.from_bytes(decompressed[i:i+4], 'little')
                    tag_id = record_header & 0x3FF
                    size = (record_header >> 20) & 0xFFF
                    
                    if size == 0xFFF:
                        if i + 8 > len(decompressed):
                            break
                        size = int.from_bytes(decompressed[i+4:i+8], 'little')
                        i += 8
                    else:
                        i += 4

                    if i + size > len(decompressed):
                        break

                    payload = decompressed[i:i+size]
                    i += size

                    if tag_id == 67: # HWPTAG_PARA_TEXT 문단 태그 ID
                        para_text = payload.decode('utf-16le', errors='ignore')
                        clean_para = re.sub(r'[\x00-\x09\x0b-\x1f]', '', para_text).strip()
                        if clean_para:
                            text_list.append(clean_para)

            ole.close()

            if not text_list:
                raise FilePreprocessError("HWP 문서에서 추출 가능한 텍스트가 없습니다.")

            return "\n".join(text_list)

        except FilePreprocessError:
            raise
        except Exception as e:
            raise FilePreprocessError(f"HWP 파싱 실패 (암호화되었거나 손상된 파일일 수 있습니다: {str(e)})")


# =========================================================
# [Step 2] 로컬 AI(Ollama) 메타데이터 생성 클래스 (FileAnalyzer)
# =========================================================
class FileAnalyzer:
    """추출된 원문을 기반으로 로컬 LLM/Vision 모델에 요청하여 메타데이터 JSON을 만드는 클래스"""

    def __init__(
        self, 
        ollama_url: str = "http://localhost:11434", 
        text_model: str = "gemma2:9b",
        vision_model: str = "llava"
    ):
        """Ollama API URL 및 사용할 텍스트/비전 LLM 모델명 초기화"""
        self.ollama_url = ollama_url.rstrip("/")
        self.text_model = text_model
        self.vision_model = vision_model

    def _get_file_info(self, file_path: str) -> Dict[str, Any]:
        """파일의 원본 이름, 확장자, 바이트 크기 및 분석 시각 메타데이터 구성"""
        file_name = os.path.basename(file_path)
        ext = os.path.splitext(file_name)[1].lower()
        size = os.path.getsize(file_path) if os.path.exists(file_path) else 0
        
        return {
            "original_name": file_name,
            "file_extension": ext,
            "file_size_bytes": size,
            "analyzed_at": datetime.now().isoformat()
        }

    # ---------------------------------------------------------
    # 텍스트 기반 문서 분석 및 메타데이터 JSON 생성
    # ---------------------------------------------------------
    def analyze_document_text(self, file_path: str, extracted_text: str) -> Dict[str, Any]:
        """문서 텍스트 원문을 로컬 텍스트 LLM에 전달하여 제목, 태그, 요약 JSON을 추출하는 함수"""
        file_name = os.path.basename(file_path)
        file_info = self._get_file_info(file_path)

        prompt = f"""
You are a professional file metadata analyzer. Analyze the provided text content and generate structured metadata in JSON format.

[File Information]
- Original Filename: {file_name}
- Content Text:
{extracted_text}

[Output Requirements]
Return ONLY a valid JSON object with the following keys:
1. "display_name": A clean, concise, and descriptive title for the file in Korean (Do NOT include file extension).
2. "tags": An array of 3 to 5 relevant keyword strings (without '#' symbol).
3. "description": A brief 1-2 sentence summary of the content in Korean.

Example JSON output format:
{{
  "display_name": "JSL 일본어 초급 교재 1권",
  "tags": ["JSL", "일본어", "초급", "교재"],
  "description": "JSL 일본어 초급 학습용 문법 및 단어 교재입니다."
}}
"""

        payload = {
            "model": self.text_model,
            "prompt": prompt,
            "format": "json", # Ollama에 JSON 출력 형식 강제 설정
            "stream": False,
            "options": {
                "temperature": 0.2,
                "num_predict": 400
            }
        }

        try:
            # Ollama 서버 API 호출
            response = OllamaManager.request("generate", payload, timeout=120, base_url=self.ollama_url)
            response.raise_for_status()

            res_data = response.json()
            raw_response_text = res_data.get("response", "").strip()

            # 정규표현식으로 순수 JSON 영역만 파싱
            match = re.search(r'\{.*\}', raw_response_text, re.DOTALL)
            json_str = match.group(0) if match else raw_response_text
            parsed_json = json.loads(json_str)

            tags = parsed_json.get("tags", [])
            desc = parsed_json.get("description", "")
            tags_formatted = ", ".join([f"#{t}" for t in tags]) if tags else "#일반"
            ai_comment_str = f"태그: {tags_formatted} / 코멘트: {desc}"

            # ✨ [@TYPE: @DB 추가] MainProcessor 순서도 라우팅과 규격 통일
            return {
                "@TYPE": "@DB",
                "status": "SUCCESS",
                "file_info": file_info,
                "metadata": {
                    "@TYPE": "@DB",
                    "display_name": parsed_json.get("display_name", file_info["original_name"].rsplit('.', 1)[0]),
                    "tags": tags,
                    "description": desc,
                    "ai_comment": ai_comment_str, # GUI 디스플레이용 완성형 문자열
                    "ocr_text": ""
                },
                "error": None
            }

        except requests.exceptions.ConnectionError:
            return self._build_fallback_response(file_info, "Ollama AI 서버에 연결할 수 없습니다. (Ollama 실행 필요)")
        except requests.exceptions.Timeout:
            return self._build_fallback_response(file_info, "AI 분석 시간 초과 (Timeout - 응답 지연)")
        except json.JSONDecodeError:
            return self._build_fallback_response(file_info, "AI 응답 파싱 실패 (유효하지 않은 JSON 구조)")
        except Exception as e:
            return self._build_fallback_response(file_info, f"Text AI 분석 중 예외 발생 ({str(e)})")

    # ---------------------------------------------------------
    # 이미지 바이너리 기반 Vision 멀티모달 분석
    # ---------------------------------------------------------
    def analyze_image_bytes(self, file_path: str, img_bytes: bytes) -> Dict[str, Any]:
        """이미지 바이트 데이터를 Base64로 인코딩하여 Vision 모델(llava)에 전달하는 함수"""
        import base64
        
        file_name = os.path.basename(file_path)
        file_info = self._get_file_info(file_path)

        # 이미지를 Base64 문자열로 변환
        base64_img = base64.b64encode(img_bytes).decode('utf-8')

        prompt = f"""
You are an image analysis expert. Analyze the provided image and return a JSON object with:
1. "display_name": A descriptive name for this image in Korean (without extension).
2. "tags": An array of 3 to 5 relevant keyword strings in Korean (without '#' symbol).
3. "description": A brief 1-2 sentence description/summary of what is shown in the image in Korean.

Filename: {file_name}

Example JSON output format:
{{
  "display_name": "팀 회의 화이트보드 메모",
  "tags": ["회의", "아이디어", "아키텍처", "일정"],
  "description": "팀 프로젝트 아키텍처 및 세부 일정이 기록된 화이트보드 이미지입니다."
}}
"""

        payload = {
            "model": self.vision_model,
            "messages": [
                {
                    "role": "user",
                    "content": prompt,
                    "images": [base64_img]
                }
            ],
            "format": "json",
            "stream": False,
            "options": {
                "temperature": 0.2,
                "num_predict": 200
            },
            "keep_alive": "5m"
        }

        try:
            response = OllamaManager.request("chat", payload, timeout=180, base_url=self.ollama_url)
            response.raise_for_status()

            res_data = response.json()
            raw_response_text = res_data.get("message", {}).get("content", "").strip()

            match = re.search(r'\{.*\}', raw_response_text, re.DOTALL)
            json_str = match.group(0) if match else raw_response_text
            parsed_json = json.loads(json_str)

            tags = parsed_json.get("tags", [])
            desc = parsed_json.get("description", "")
            tags_formatted = ", ".join([f"#{t}" for t in tags]) if tags else "#이미지"
            ai_comment_str = f"태그: {tags_formatted} / 코멘트: {desc}"

            # ✨ [@TYPE: @DB 추가] MainProcessor 순서도 라우팅과 규격 통일
            return {
                "@TYPE": "@DB",
                "status": "SUCCESS",
                "file_info": file_info,
                "metadata": {
                    "@TYPE": "@DB",
                    "display_name": parsed_json.get("display_name", file_info["original_name"].rsplit('.', 1)[0]),
                    "tags": tags,
                    "description": desc,
                    "ai_comment": ai_comment_str, # GUI 디스플레이용 완성형 문자열
                    "ocr_text": ""
                },
                "error": None
            }

        except requests.exceptions.ConnectionError:
            return self._build_fallback_response(file_info, "Ollama AI 서버에 연결할 수 없습니다. (Ollama 실행 필요)")
        except requests.exceptions.Timeout:
            return self._build_fallback_response(file_info, "Vision AI 분석 시간 초과 (Timeout - 로딩 또는 연산 지연)")
        except json.JSONDecodeError:
            return self._build_fallback_response(file_info, "Vision AI 응답 파싱 실패 (유효하지 않은 JSON 구조)")
        except Exception as e:
            return self._build_fallback_response(file_info, f"Vision AI 분석 중 예외 발생 ({str(e)})")

    # ---------------------------------------------------------
    # 예외 발생 시 안전하게 기본값을 채워주는 폴백(Fallback) 함수
    # ---------------------------------------------------------
    def _build_fallback_response(self, file_info: Dict[str, Any], error_message: str) -> Dict[str, Any]:
        """AI 분석 실패 시 프로그램 멈춤 없이 최소 메타데이터로 구성된 실패 응답 JSON 반환"""
        default_name = file_info["original_name"].rsplit('.', 1)[0]
        return {
            "@TYPE": "@DB",
            "status": "FAILED",
            "file_info": file_info,
            "metadata": {
                "@TYPE": "@DB",
                "display_name": default_name,
                "tags": [],
                "description": f"분석 실패: {error_message}",
                "ai_comment": f"#분석실패 / 코멘트: {error_message}",
                "ocr_text": ""
            },
            "error": error_message
        }


# =====================================================================
# 통합 단독 테스트 실행부 (main)
# =====================================================================
if __name__ == "__main__":
    # 파서 및 분석기 객체 인스턴스 생성
    extractor = TextExtractor(max_chars=2000, max_img_size=512)
    analyzer = FileAnalyzer(text_model="qwen2.5:3b", vision_model="llava")

    # 테스트할 파일 샘플 리스트
    test_files = [
        "sample.txt",
        "sample.pdf",
        "sample.docx",
        "sample.xlsx",
        "sample.pptx",
        "sample.hwp",
        "sample.hwpx",
        "sample.json",
        "sample.png",
        "sample.gif",
        "sample.mp3",
        "sample.mp4",
        "sample.zip"  # 압축 파일 거부 차단 동작 테스트
    ]

    print("=== [통합 테스트] 문서/이미지/미디어 파일 전처리 및 AI 분석 ===")

    for file_path in test_files:
        print(f"\n========================================")
        print(f"📄 대상 파일: {file_path}")

        if not os.path.exists(file_path):
            print(f"[경고] '{file_path}' 파일을 찾을 수 없습니다. 테스트를 스킵합니다.")
            continue

        # 🖼️ 1. 이미지 파일인 경우의 분기 처리
        if extractor.is_image_file(file_path):
            img_bytes, status = extractor.process_image(file_path)
            print(f"[Step 1 이미지 전처리 상태]: {status}")

            if status == "SUCCESS":
                print(f"[Step 2 Vision AI({analyzer.vision_model}) 분석 요청 중...]")
                ai_result = analyzer.analyze_image_bytes(file_path, img_bytes)
                print("[Step 2 AI 분석 결과 (JSON)]:")
                print(json.dumps(ai_result, ensure_ascii=False, indent=2))
            else:
                print(f"[Step 1 예외/거부 사유]: {status}")

        # 🎵 2. 오디오/비디오 미디어 파일인 경우의 분기 처리
        elif extractor.is_media_file(file_path):
            extracted_text, status = extractor.process_media(file_path)
            print(f"[Step 1 음성 인식(STT) 상태]: {status}")

            if status == "SUCCESS":
                print(f"[Step 1 추출 텍스트 (최대 300자)]:\n---")
                print(extracted_text[:300] + ("..." if len(extracted_text) > 300 else ""))
                print("---")

                print(f"[Step 2 Text AI({analyzer.text_model}) 분석 요청 중...]")
                ai_result = analyzer.analyze_document_text(file_path, extracted_text)
                print("[Step 2 AI 분석 결과 (JSON)]:")
                print(json.dumps(ai_result, ensure_ascii=False, indent=2))
            else:
                print(f"[Step 1 예외/거부 사유]: {status}")

        # 📝 3. 일반 문서 파일인 경우의 분기 처리
        else:
            extracted_text, status = extractor.extract(file_path)
            print(f"[Step 1 텍스트 추출 상태]: {status}")

            if status == "SUCCESS":
                print(f"[Step 1 추출 텍스트 (최대 300자)]:\n---")
                print(extracted_text[:300] + ("..." if len(extracted_text) > 300 else ""))
                print("---")

                print(f"[Step 2 Text AI({analyzer.text_model}) 분석 요청 중...]")
                ai_result = analyzer.analyze_document_text(file_path, extracted_text)
                print("[Step 2 AI 분석 결과 (JSON)]:")
                print(json.dumps(ai_result, ensure_ascii=False, indent=2))
            else:
                print(f"[Step 1 예외/거부 사유]: {status}")
#========================file_pipeline.py=====================
#========================main_processor.py=====================

# =========================================================
# [main_processor.py]
# 통합 코어 모듈 래퍼 - 하위 호환성 유지
# =========================================================
# 이 파일은 하위 호환성을 위해 유지되며,
# 새로운 코드에서는 core.py의 ClasqCore를 직접 사용하세요.
from .core import ClasqCore as MainProcessor

#========================main_processor.py=====================
#========================query_parser.py=====================

# =========================================================
# [query_parser.py] 
# 프론트엔드 자연어 입력 파싱 및 의도 분류 모듈 (@검색, @대화)
# =========================================================
import re          # 정규표현식 모듈
import json        # JSON 데이터 디코딩/인코딩 모듈
from ollama_manager import OllamaManager
from typing import Dict, Any


class SearchQueryParser:
    """
    [자연어 분석 모듈] 
    사용자가 입력한 검색어/일상대화 문장을 로컬 AI(Ollama)에 전달하고,
    의도를 분석하여 구조화된 JSON(@검색, @대화)으로 변환해 반환합니다.
    """

    def __init__(self, ollama_url: str = "http://localhost:11434", model: str = "gemma2:9b"):
        self.ollama_url = ollama_url.rstrip("/")
        self.model = model

    def parse_user_query(self, user_text: str) -> Dict[str, Any]:
        """사용자 입력 자연어를 분석하여 '@TYPE'이 포함된 JSON 객체 반환"""
        
        prompt = f"""
You are a smart Assistant for a File Management System.
Analyze the user's input string and classify the intent into either '@검색' or '@대화'.

User Input: "{user_text}"

[Classification Rules]
1. Set "@TYPE" to "@검색" IF:
   - The user wants to find, search, show, or list local files/documents/images.
   - Examples: "pdf 파일 찾아줘", "지난주 회의록 어디 있어?", "jpg 이미지 보여줘"
   - Extract key search terms into "query_keywords" (array of strings).
   - Extract file extensions if explicitly mentioned into "target_extension" (e.g., [".pdf"], [".xlsx"]).

2. Set "@TYPE" to "@대화" IF:
   - The user is making casual greetings, small talk, or general questions NOT related to searching local files.
   - Examples: "안녕", "오늘 날씨 어때?", "넌 누구야?"
   - You MUST generate a polite, complete, and helpful Korean response in "reply_text".
   - DO NOT just echo or repeat the user's input! Provide a helpful real answer.

[Output Format Requirements]
Return ONLY a valid JSON object matching one of these structures:

If "@검색":
{{
  "@TYPE": "@검색",
  "query_keywords": ["keyword1", "keyword2"],
  "target_extension": [".pdf"],
  "raw_query": "{user_text}"
}}

If "@대화":
{{
  "@TYPE": "@대화",
  "reply_text": "사용자 질문에 맞는 친절하고 완성도 높은 한글 대화 응답 문장"
}}

Keep your response brief and concise.
"""

        payload = {
            "model": self.model,
            "prompt": prompt,
            "format": "json",
            "stream": False,
            "options": {"temperature": 0.1, "num_predict": 100}
        }

        try:
            res = OllamaManager.request("generate", payload, timeout=120, base_url=self.ollama_url)
            res.raise_for_status()
            
            raw_text = res.json().get("response", "").strip()

            # 응답 내 Pure JSON 영역 추출
            match = re.search(r'\{.*\}', raw_text, re.DOTALL)
            json_str = match.group(0) if match else raw_text
            parsed_json = json.loads(json_str)

            # @TYPE 누락 시 기본 폴백
            if "@TYPE" not in parsed_json:
                parsed_json["@TYPE"] = "@대화"
                parsed_json["reply_text"] = raw_text if raw_text else "안녕하세요! 무엇을 도와드릴까요?"

            return {
                "status": "SUCCESS", 
                "data": parsed_json, 
                "error": None
            }

        except Exception as e:
            return {
                "status": "FAILED",
                "data": {
                    "@TYPE": "@ERROR", 
                    "message": f"자연어 파싱 처리 중 오류 발생: {str(e)}"
                },
                "error": str(e)
            }


# =========================================================
# 단독 테스트 실행부 (main)
# =========================================================
if __name__ == "__main__":
    parser = SearchQueryParser(model="gemma2:9b")

    print("=== [SearchQueryParser] 자연어 의도 파싱 테스트 ===")

    res1 = parser.parse_user_query("지난주에 만든 프로젝트 보고서 pdf 파일 찾아줘")
    print("\n[테스트 1 - 검색 요청 결과]:\n", json.dumps(res1, ensure_ascii=False, indent=2))

    res2 = parser.parse_user_query("안녕, 너는 어떤 일을 할 수 있니?")
    print("\n[테스트 2 - 대화 요청 결과]:\n", json.dumps(res2, ensure_ascii=False, indent=2))

#========================query_parser.py=====================
#========================search_engine.py=====================
# =========================================================
# [search_engine.py]
# DB 검색 및 자연어 의도 라우팅 후속 로직 처리 모듈
# (불용어 제거, 동의어 사전 확장, 0건 방지 폴백 검색 완결판)
# =========================================================
import sqlite3
from typing import Dict, Any, List


class SearchEngine:
    """
    [핵심 후속 처리 엔진]
    query_parser 및 main_processor에서 넘겨받은 JSON 데이터('@TYPE')를 확인하여
    1) DB(files 테이블) 조건 조회 및 결과 테이블 반환(@검색)
    2) AI 대화 메시지 팝업 전달(@대화)
    의 실제 후속 액션을 담당하는 클래스입니다.
    """

    # 1. 자연어 검색 품질 향상을 위한 확장 불용어(Stopwords) 세트
    STOP_WORDS = {
        "파일", "문서", "폴더", "데이터", "자료", "내용", "것",
        "찾아줘", "보여줘", "검색", "알려줘", "꺼내줘", "어디있어", "어디", "있냐",
        "관련된", "관련", "에", "대한", "중 중에서", "중", "내", "속", "제일", "최근", "좀", "하나",
        "pdf", "hwp", "hwpx", "docx", "xlsx", "pptx", "png", "jpg", "jpeg", "gif", "mp3", "mp4"
    }

    # 2. 검색 정확도 극대화를 위한 동의어/유의어 매핑 사전 (🌟 '전쟁' 유의어 추가)
    SYNONYM_MAP = {
        "실습": ["실습", "현장실습", "인턴", "교육"],
        "학교": ["학교", "캠퍼스", "학사"],
        "노래": ["노래", "음원", "가사", "음악", "작업"],
        "번안": ["번안", "번역", "가사"],
        "이미지": ["이미지", "사진", "그림", "gif", "png", "jpg"],
        "보고서": ["보고서", "리포트", "과제", "기안서"],
        "회의": ["회의", "미팅", "회의록"],
        "전쟁": ["전쟁", "대전", "전투"]
    }

    def __init__(self, db_path: str = "file_manager.db"):
        """검색에 사용할 SQLite 데이터베이스 파일 경로 초기화"""
        self.db_path = db_path

    # ---------------------------------------------------------
    # [1] 자연어 파싱 결과 분기 및 액션 제어 함수
    # ---------------------------------------------------------
    def process_query_result(self, parsed_data: Dict[str, Any]) -> Dict[str, Any]:
        """
        '@TYPE'값(@검색, @대화, @ERROR)에 따라 UI가 실행할 액션 명령 포장
        """
        type_val = parsed_data.get("@TYPE")

        # [Case 1] 검색 -> DB 조회 후 표 데이터 갱신 명령
        if type_val in ["search", "@검색"]:

            condition = parsed_data.get("condition", {})
            if condition:
                raw_keywords = condition.get("tags", [])
            else:
                raw_keywords = parsed_data.get("query_keywords", [])

            exts = parsed_data.get("target_extension", [])

            split_keywords = []
            for kw in raw_keywords:
                split_keywords.extend(kw.split())  # 띄어쓰기 기준으로 단어 분리

            # 불용어 제거 필터링
            filtered_keywords = [
                kw.strip().lower() for kw in split_keywords
                if kw.strip() and kw.strip().lower() not in self.STOP_WORDS
            ]

            final_keywords = filtered_keywords if filtered_keywords else [
                kw.strip() for kw in split_keywords if kw.strip()]

            # 1차: 엄격한 검색 (AND) -> 안되면 2차: 완화된 검색 (OR)
            search_results, is_fallback = self.search_files_smart(
                final_keywords, exts)

            display_kw = ', '.join(final_keywords) if final_keywords else "전체"

            if is_fallback:
                msg = f"'{display_kw}' 완벽 일치 항목이 없어 연관 키워드 검색 결과 {len(search_results)}건을 보여드립니다."
            else:
                msg = f"'{display_kw}' 검색 결과 {len(search_results)}건을 찾았습니다."

            return {
                "action": "UPDATE_TABLE",
                "message": msg,
                "data": search_results
            }

        # [Case 2] @대화 -> AI 대화 응답 출력 명령
        elif type_val == "@대화":
            reply = parsed_data.get("reply_text", "안녕하세요! 무엇을 도와드릴까요?")
            return {
                "action": "SHOW_CHAT",
                "message": reply,
                "data": []
            }

        # Ⓒ [Case 3] 오류 및 예외
        else:
            return {
                "action": "ERROR",
                "message": parsed_data.get("message", "알 수 없거나 올바르지 않은 요청 타입입니다."),
                "data": []
            }

    # ---------------------------------------------------------
    # [2] 지능형 DB 검색 및 Fallback 제어 로직
    # ---------------------------------------------------------
    def search_files_smart(self, keywords: List[str], exts: List[str] = None) -> tuple[List[tuple], bool]:
        """
        1차(AND 검색) 시도 후 결과가 0건이면 2차(OR 완화 검색)로 자동 전환
        :return: (검색결과 리스트, Fallback 적용 여부)
        """
        if not keywords and not exts:
            return self._execute_sql_query([], exts, match_mode="AND"), False

        # 1차 시도: 동의어 적용 AND 조건 검색
        results = self._execute_sql_query(keywords, exts, match_mode="AND")
        if results:
            return results, False

        # 2차 시도 (Fallback): 1차에서 0건이면 OR 조건으로 완화 검색
        results_or = self._execute_sql_query(keywords, exts, match_mode="OR")
        return results_or, True

    def _execute_sql_query(self, keywords: List[str], exts: List[str] = None, match_mode: str = "AND") -> List[tuple]:
        """실제 SQLite LIKE SQL 문을 생성하고 실행하는 내부 함수"""
        conn = sqlite3.connect(self.db_path)
        cursor = conn.cursor()

        query = "SELECT id, file_name, file_path, ai_comment, category, tags FROM files WHERE 1=1"
        params = []

        if keywords:
            keyword_group_sql = []

            for kw in keywords:
                if not kw.strip():
                    continue

                # 동의어 사전 매핑을 통한 검색어 확장
                synonyms = self.SYNONYM_MAP.get(kw, [kw])

                # 각 단어 또는 동의어 그룹 내에서 OR 매칭 조건 형성
                synonym_conditions = []
                for syn in synonyms:
                    synonym_conditions.append(
                        "(file_name LIKE ? OR ai_comment LIKE ? OR category LIKE ? OR tags LIKE ?)")
                    params.extend([f"%{syn}%", f"%{syn}%", f"%{syn}%", f"%{syn}%"])

                single_kw_sql = "(" + " OR ".join(synonym_conditions) + ")"
                keyword_group_sql.append(single_kw_sql)

            if keyword_group_sql:
                # AND 모드와 OR 모드 분기 (match_mode가 OR일 경우 하나만 걸려도 매칭되도록 완화)
                join_operator = " AND " if match_mode == "AND" else " OR "
                query += " AND (" + join_operator.join(keyword_group_sql) + ")"

        # 확장자 필터 (예: .pdf, .docx 등)
        if exts:
            ext_conditions = []
            for ext in exts:
                if ext.strip():
                    ext_conditions.append("file_path LIKE ?")
                    params.append(f"%{ext}")

            if ext_conditions:
                query += " AND (" + " OR ".join(ext_conditions) + ")"

        cursor.execute(query, params)
        results = cursor.fetchall()
        conn.close()

        return results


# =========================================================
# 단독 테스트 실행부 (main)
# =========================================================
if __name__ == "__main__":
    search_engine = SearchEngine(db_path="file_manager.db")

    print("=== [SearchEngine] 스마트 검색 및 동의어 테스트 ===")

    sample_search_json = {
        "@TYPE": "@검색",
        "query_keywords": ["전쟁", "파일"],
        "target_extension": []
    }
    print("\n[검색 파싱 결과 처리]:\n",
          search_engine.process_query_result(sample_search_json))
#========================search_engine.py=====================
#========================workers.py=====================
import os
from PySide6.QtCore import QThread, Signal
from .core import ClasqCore
from .query_parser import SearchQueryParser


class FolderScanAndTagWorker(QThread):
    progress = Signal(str)
    finished = Signal()
    error = Signal(str)

    def __init__(self, folder_paths: list, core: ClasqCore):
        super().__init__()
        self.folder_paths = folder_paths
        self.core = core

    # 비동기 제어랑 스레드 반복 처리 하는 함수
    # 독립된 백그라운드 스레드인 QThread 안에서 동작함
    def run(self):
        try:
            valid_extensions = (
                '.txt', '.pdf', '.docx', '.xlsx', '.pptx', '.hwp', '.hwpx',
                '.jpg', '.jpeg', '.png', '.webp', '.bmp', '.gif',
                '.mp3', '.mp4', '.wav', '.m4a', '.mkv', '.avi'
            )
            files_to_process = []

            for target_path in self.folder_paths:
                clean_target = os.path.abspath(os.path.normpath(target_path))
                if os.path.isfile(clean_target):
                    if clean_target.lower().endswith(valid_extensions):
                        files_to_process.append(clean_target)
                    continue
                if not os.path.isdir(clean_target):
                    continue
                for root, _, files in os.walk(clean_target):
                    for file in files:
                        if file.lower().endswith(valid_extensions):
                            full_path = os.path.join(root, file)
                            files_to_process.append(os.path.abspath(os.path.normpath(full_path)))

            files_to_process = list(dict.fromkeys(files_to_process))
            if not files_to_process:
                self.error.emit("스캔할 지원 파일이 지정된 경로에 없습니다.")
                return

            total_count = len(files_to_process)

            for idx, file_path in enumerate(files_to_process, start=1):
                # 파일 리스트를 하나씩 꺼내서
                # core.process_file_upload로 넘겨줌
                file_name = os.path.basename(file_path)
                self.progress.emit(
                    # 파일 처리 진행 상황 프론트 엔드 연결 필요
                    f"AI 분석 중 ({idx}/{total_count}): {file_name}")
                result = self.core.process_file_upload(file_path)
                db_save = result.get("db_save", {})
                if not db_save.get("success"):
                    detail = db_save.get("message", "원인을 확인할 수 없습니다.")
                    raise RuntimeError(f"'{file_name}'의 DB 저장에 실패했습니다. {detail}")

            self.finished.emit()

        except Exception as e:
            self.error.emit(f"스캔 및 태깅 작업 중 오류 발생: {str(e)}")


class QueryParseWorker(QThread):
    # 사용자가 자연어 입력 시 UI 멈춤 방지
    finished = Signal(dict)
    error = Signal(str)

    def __init__(self, user_text: str, query_parser: SearchQueryParser):
        super().__init__()
        self.user_text = user_text
        self.query_parser = query_parser

    def run(self):
        try:
            result = self.query_parser.parse_user_query(self.user_text)
            self.finished.emit(result)
        except Exception as e:
            self.error.emit(f"자연어 파싱 처리 중 오류: {str(e)}")

#========================workers.py=====================