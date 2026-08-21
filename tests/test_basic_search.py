import os
import shutil
import sqlite3
import time
import unittest
from pathlib import Path
from unittest import mock

from pptx import Presentation
from docx import Document

from src.ui.views.search_view import SearchView
from src.utils.db_manager import FileRegistryManager
from src.utils.local_text_index import LocalTextIndexer
from src.utils.search_engine import SearchEngine
from src.utils.search_aliases import build_search_alias_map
from src.utils.search_normalization import strip_korean_particle


class CountingPptxExtractor:
    def __init__(self):
        self.calls = 0

    def _read_pptx(self, path):
        self.calls += 1
        presentation = Presentation(path)
        return "\n".join(
            shape.text.strip()
            for slide in presentation.slides
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        )


class CountingIndexExtractor:
    def __init__(self, texts):
        self.texts = texts
        self.calls = 0

    def extract_for_index(self, path):
        self.calls += 1
        return self.texts[Path(path).suffix.lower()], "SUCCESS"


class StatusIndexExtractor:
    def __init__(self, status):
        self.status = status

    def extract_for_index(self, _path):
        return ("searchable truncated text" if self.status == "TRUNCATED" else ""), self.status


class BasicSearchTests(unittest.TestCase):
    def setUp(self):
        self.root = Path("tests/fixtures/basic_search_runtime").resolve()
        self.root.mkdir(parents=True, exist_ok=True)
        self.db_path = str(self.root / "search.db")
        self.registry = FileRegistryManager(self.db_path)

    def tearDown(self):
        shutil.rmtree(self.root, ignore_errors=True)

    def _pptx(self, name="random_1234.pptx", text="중소기업 경영 컨설팅 결과 보고서"):
        path = self.root / name
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        textbox = slide.shapes.add_textbox(0, 0, 5000000, 1000000)
        textbox.text = text
        presentation.save(path)
        return str(path)

    def _cache(self, path):
        stat = os.stat(path)
        self.registry.cache_file_fingerprints([{
            "file_path": path,
            "file_hash": self.registry.compute_file_hash(path),
            "file_size": stat.st_size,
            "file_mtime_ns": stat.st_mtime_ns,
        }])

    def test_rules_parser_normalizes_presentation_aliases(self):
        expected = {"query_keywords": ["컨설팅"], "target_extension": ["ppt", "pptx"]}
        for query in ("컨설팅 관련 피피티 찾아줘", "컨설팅 PPT 찾아줘", "컨설팅 파워포인트"):
            parsed = SearchView._parse_natural_query(None, query)
            self.assertEqual(parsed["query_keywords"], expected["query_keywords"])
            self.assertEqual(parsed["target_extension"], expected["target_extension"])
        parsed = SearchView._parse_natural_query(None, "PPTX 찾아줘")
        self.assertEqual(parsed["query_keywords"], [])
        self.assertEqual(parsed["target_extension"], ["pptx"])
        self.assertEqual(
            SearchView._parse_natural_query(None, "발표자료 찾아줘")["query_keywords"],
            ["발표자료"],
        )

    def test_rules_parser_strips_particles_and_expands_file_types(self):
        cases = {
            "회의록을 찾아줘": (["회의록"], []),
            "보고서를 보여줘": (["보고서"], []),
            "PDF 문서 보여줘": ([], ["pdf"]),
            "엑셀 파일 찾아줘": ([], ["xlsx"]),
            "마크다운 파일 찾아줘": ([], ["md", "markdown"]),
        }
        for query, expected in cases.items():
            parsed = SearchView._parse_natural_query(None, query)
            self.assertEqual((parsed["query_keywords"], parsed["target_extension"]), expected)

    def test_particle_stripping_is_conservative(self):
        expected = {
            "회의록을": "회의록", "보고서에서": "보고서",
            "게이트웨이": "게이트웨이",
            "국가": "국가", "세로": "세로", "카테고리": "카테고리",
        }
        for token, normalized in expected.items():
            self.assertEqual(strip_korean_particle(token), normalized)

    def test_request_phrases_are_removed_after_particle_normalization(self):
        parsed = SearchView._parse_natural_query(
            None, "창업 폴더에 있는 로보독 발표자료 찾아줘"
        )
        self.assertEqual(parsed["query_keywords"], ["창업", "로보독", "발표자료"])

    def test_configurable_aliases_are_one_keyword_group(self):
        paths = []
        for name, body in (("SmartCare_panel.txt", "patient display"),
                           ("LoboDoc_plan.txt", "business plan"),
                           ("network_switch.txt", "MAC 주소 테이블")):
            path = self.root / name
            path.write_text(body, encoding="utf-8")
            self._cache(str(path))
            paths.append(str(path))
        LocalTextIndexer(self.db_path).synchronize(paths)
        engine = SearchEngine(self.db_path, project_aliases={
            "스마트케어": ("smartcare",),
            "smartcare": ("스마트케어",),
            "로보독": ("lobodoc",),
            "lobodoc": ("로보독",),
        })
        rows, _ = engine.search_files_smart(["스마트케어"])
        self.assertEqual(rows[0][2], paths[0])
        self.assertEqual(engine.get_result_metadata(paths[0])["keyword_count"], 1)
        rows, _ = engine.search_files_smart(["로보독"])
        self.assertEqual(rows[0][2], paths[1])
        self.assertEqual(engine.get_result_metadata(paths[1])["keyword_count"], 1)
        rows, _ = engine.search_files_smart(["맥", "주소", "테이블"])
        self.assertEqual(rows[0][2], paths[2])

    def test_product_defaults_exclude_project_aliases(self):
        with mock.patch.dict(os.environ, {"SEARCH_PROJECT_ALIASES_JSON": ""}):
            aliases = build_search_alias_map()
        self.assertIn("mac", aliases["맥"])
        self.assertNotIn("스마트케어", aliases)
        self.assertNotIn("로보독", aliases)

    def test_compact_english_variants_match(self):
        path = self.root / "smart-care_lobo-doc.txt"
        path.write_text("fixture", encoding="utf-8")
        self._cache(str(path))
        engine = SearchEngine(self.db_path)
        self.assertEqual(engine.search_files_smart(["smartcare"])[0][0][2], str(path))
        self.assertEqual(engine.search_files_smart(["lobodoc"])[0][0][2], str(path))

    def test_hangul_is_only_a_strong_alias_with_file_context(self):
        self.assertEqual(
            SearchView._parse_natural_query(None, "한글 파일 찾아줘")["target_extension"],
            ["hwp", "hwpx"],
        )
        parsed = SearchView._parse_natural_query(None, "한글 OCR 자료")
        self.assertEqual(parsed["target_extension"], [])
        self.assertIn("한글", parsed["query_keywords"])

    def test_docx_extractor_and_local_body_search(self):
        path = self.root / "random_name.docx"
        document = Document()
        document.add_paragraph("중소기업 컨설팅 결과 보고서")
        document.save(path)
        self._cache(str(path))

        stats = LocalTextIndexer(self.db_path).synchronize([str(path)])
        self.assertEqual(stats["success"], 1)
        parsed = SearchView._parse_natural_query(None, "컨설팅 워드 찾아줘")
        rows = SearchEngine(self.db_path).process_query_result(parsed)["data"]
        self.assertEqual([row[2] for row in rows], [str(path)])

    def test_multi_format_body_search_and_incremental_indexing(self):
        paths = []
        texts = {
            ".pdf": "사업계획 발표 자료",
            ".md": "스마트폴더 검색 정확도 개선",
            ".txt": "회의록 검색 테스트",
        }
        for extension in texts:
            path = self.root / f"random_name{extension}"
            path.write_text("fixture", encoding="utf-8")
            self._cache(str(path))
            paths.append(str(path))
        extractor = CountingIndexExtractor(texts)
        indexer = LocalTextIndexer(self.db_path, extractor=extractor)
        first = indexer.synchronize(paths)
        second = indexer.synchronize(paths)
        self.assertEqual(first["success"], 3)
        self.assertEqual(second["unchanged"], 3)
        self.assertEqual(extractor.calls, 3)

        engine = SearchEngine(self.db_path)
        for query, expected_suffix in (
            ("사업계획 PDF 찾아줘", ".pdf"),
            ("검색 정확도 마크다운 찾아줘", ".md"),
        ):
            rows = engine.process_query_result(
                SearchView._parse_natural_query(None, query)
            )["data"]
            self.assertEqual(Path(rows[0][2]).suffix, expected_suffix)

        changed_path = Path(paths[0])
        changed_path.write_text("changed", encoding="utf-8")
        future_ns = max(os.stat(changed_path).st_mtime_ns + 1_000_000, time.time_ns())
        os.utime(changed_path, ns=(future_ns, future_ns))
        changed = indexer.synchronize(paths)
        self.assertEqual(changed["indexed"], 1)
        self.assertEqual(extractor.calls, 4)

        changed_path.unlink()
        deleted = indexer.synchronize(paths[1:])
        self.assertEqual(deleted["deleted"], 1)

    def test_separator_normalization_and_multi_field_evidence(self):
        path = self.root / "LoboDoc_사업계획_v3.txt"
        path.write_text("LoboDoc 사업계획", encoding="utf-8")
        self._cache(str(path))
        LocalTextIndexer(self.db_path).synchronize([str(path)])
        engine = SearchEngine(self.db_path)
        parsed = SearchView._parse_natural_query(None, "LoboDoc 사업계획 파일 찾아줘")
        rows = engine.process_query_result(parsed)["data"]
        self.assertEqual(rows[0][2], str(path))
        metadata = engine.get_result_metadata(str(path))
        self.assertEqual(metadata["keyword_matches"], 2)
        self.assertIn("filename", metadata["match_source"])
        self.assertIn("text", metadata["match_source"])
        breakdown = metadata["score_breakdown"]
        self.assertGreater(breakdown["filename_score"], 0)
        self.assertGreater(breakdown["text_score"], 0)
        self.assertGreater(breakdown["evidence_bonus"], 0)
        self.assertEqual(breakdown["keyword_coverage"], 1.0)
        self.assertEqual(breakdown["final_score"], metadata["relevance_score"])

    def test_no_text_and_truncated_index_statuses_are_distinct(self):
        no_text_path = self.root / "scan.pdf"
        no_text_path.write_bytes(b"fixture")
        truncated_path = self.root / "large.md"
        truncated_path.write_text("fixture", encoding="utf-8")
        self._cache(str(no_text_path))
        self._cache(str(truncated_path))

        no_text = LocalTextIndexer(
            self.db_path, extractor=StatusIndexExtractor("NO_TEXT")
        ).synchronize([str(no_text_path)])
        truncated = LocalTextIndexer(
            self.db_path, extractor=StatusIndexExtractor("TRUNCATED")
        ).synchronize([str(no_text_path), str(truncated_path)])
        self.assertEqual(no_text["no_text"], 1)
        self.assertEqual(truncated["truncated"], 1)

        conn = sqlite3.connect(self.db_path)
        try:
            statuses = dict(conn.execute(
                "SELECT file_path, extract_status FROM file_text_index"
            ).fetchall())
        finally:
            conn.close()
        self.assertEqual(statuses[str(no_text_path)], "no_text")
        self.assertEqual(statuses[str(truncated_path)], "truncated")

    def test_cache_only_pptx_is_found_by_local_slide_text(self):
        path = self._pptx()
        self._cache(path)
        stats = LocalTextIndexer(self.db_path).synchronize([path])
        self.assertEqual(stats["success"], 1)
        engine = SearchEngine(self.db_path)
        parsed = SearchView._parse_natural_query(None, "컨설팅 관련 피피티 찾아줘")
        result = engine.process_query_result(parsed)
        self.assertEqual([row[2] for row in result["data"]], [path])
        metadata = engine.get_result_metadata(path)
        self.assertEqual(metadata["analysis_status"], "pending")
        self.assertIn("text", metadata["match_source"])

    def test_path_search_and_files_row_deduplication(self):
        folder = self.root / "컨설팅"
        folder.mkdir()
        path = self._pptx(name="컨설팅/ordinary.pptx", text="unrelated")
        self._cache(path)
        conn = sqlite3.connect(self.db_path)
        try:
            conn.execute(
                "INSERT INTO files(file_name,file_path,ai_comment,category) VALUES(?,?,?,?)",
                ("ordinary.pptx", path, "분석됨", "#업무"),
            )
            conn.commit()
        finally:
            conn.close()
        engine = SearchEngine(self.db_path)
        rows, _ = engine.search_files_smart(["컨설팅"], ["pptx"])
        self.assertEqual(len(rows), 1)
        self.assertEqual(engine.get_result_metadata(path)["analysis_status"], "analyzed")
        self.assertIn("path", engine.get_result_metadata(path)["match_source"])

    def test_text_index_is_incremental_and_prunes_deleted_files(self):
        path = self._pptx()
        self._cache(path)
        extractor = CountingPptxExtractor()
        indexer = LocalTextIndexer(self.db_path, extractor=extractor)
        first = indexer.synchronize([path])
        second = indexer.synchronize([path])
        self.assertEqual(first["indexed"], 1)
        self.assertEqual(second["indexed"], 0)
        self.assertEqual(second["unchanged"], 1)
        self.assertEqual(extractor.calls, 1)

        presentation = Presentation(path)
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.add_textbox(0, 0, 5000000, 1000000).text = "변경된 본문"
        presentation.save(path)
        future_ns = max(os.stat(path).st_mtime_ns + 1_000_000, time.time_ns())
        os.utime(path, ns=(future_ns, future_ns))
        changed = indexer.synchronize([path])
        self.assertEqual(changed["indexed"], 1)
        self.assertEqual(extractor.calls, 2)

        os.unlink(path)
        deleted = indexer.synchronize([])
        self.assertEqual(deleted["deleted"], 1)
        conn = sqlite3.connect(self.db_path)
        try:
            self.assertEqual(conn.execute("SELECT count(*) FROM file_text_index").fetchone()[0], 0)
        finally:
            conn.close()

    def test_legacy_ppt_is_searchable_by_name_without_content_extraction(self):
        path = self.root / "컨설팅_구형자료.ppt"
        path.write_bytes(b"legacy-ppt-placeholder")
        stats = LocalTextIndexer(self.db_path).synchronize([str(path)])
        self.assertEqual(stats["unsupported"], 1)
        engine = SearchEngine(self.db_path)
        rows, _ = engine.search_files_smart(["컨설팅"], ["ppt", "pptx"])
        self.assertEqual([row[2] for row in rows], [str(path)])
        self.assertEqual(engine.get_result_metadata(str(path))["extract_status"], "unsupported")


if __name__ == "__main__":
    unittest.main()
